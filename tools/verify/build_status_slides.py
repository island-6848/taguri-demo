#!/usr/bin/env python3
"""検証（#000006）の現在地を、企画書 1 章の流れの上に貼った 5 枚のスライドにする。

    python3 build_status_slides.py            # pptx / html / md をまとめて生成
    python3 build_status_slides.py --check    # 収まるかどうかだけ検査する

レイアウトの記法・座標系・描画は 7 枚版（tools/proposal/build_taguri_7p.py）を
そのまま使う。同じ本文領域（9.06 x 5.72 inch・4:3）に組むので、ブラウザで見た形が
そのまま pptx になる。台紙は NaU 標準スライドマスタで、見出しだけを差し込む。
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "tools/proposal"))

from pptx import Presentation  # noqa: E402

from build_taguri_7p import (  # noqa: E402
    AREA_H, AREA_W, HTML_HEAD, bars, bullets, card, check, cycle,
    draw_shape_pptx, esc, html_shape, label, para, resolve, src, stack, table,
)

TEMPLATE = ROOT / "data/formats/NaU標準スライドマスタ_2025_01.pptx"
OUT_PPTX = ROOT / "docs/000006-verification-status.pptx"
OUT_HTML = ROOT / "docs/000006-verification-status.html"
OUT_MD = ROOT / "docs/000006-verification-status.md"

DECK = "たぐりの検証 ── 現在地（2026/08/21 時点）"

# ── 1 枚目 全体像 ────────────────────────────────────────────────
S1 = dict(
    no=1, title="たぐりの検証 ── 現在地",
    claim="流れの 4 段のうち、入口と輪は実データで裏づいた。残っているのは、あらすじの抽出と母集団である。",
    stacks=[
        stack(0, AREA_W, [
            label("検証 34 件を実施した。結果を、企画書 1 章の流れ（① → ② → ③ → ④）の上に置く"),
            cycle(
                [("① 材料を集める", "履歴とクレジットは取れた。あらすじだけが取れていない"),
                 ("② 知識を作る", "順位を決めるのは名簿。内容の傾向が件数を埋める"),
                 ("③ 一覧を見る", "5 回出した。件数と順位が決まっていない"),
                 ("④ 評価を付ける", "反応が推薦を変えるところまで確かめた")],
                center="判定済み\n26 項目",
                entry=[("履歴の入口 ── 購入確認メールを検索する",
                        "成立した。7／7 で特定でき、日付は 99.2% 一致した"),
                       ("候補の母集団 ── いま観に行ける公演を集める",
                        "未決着である。被覆が標本で 81% と 43% に割れた")],
                interlude=("記録を見返す", "①② を通らず ④ へ"),
                h=2.62, nw=2.32, nh=0.66, ring_w=6.90),
            src("内訳 ── 判定済み 26 項目／納期内に測る 13 項目（P0〜P2）／納期内は測らない 6 項目（P3）／"
                "取り下げ・統合 4 項目／凍結 5 項目。ほかに検証 021・026 で起票した 7 項目のうち、"
                "3 項目は判定が付き、1 項目を取り下げ、3 項目が測定待ちである"),
        ], gap=0.10),
    ],
)

# ── 2 枚目 ① 材料を集める ────────────────────────────────────────
S2 = dict(
    no=2, title="① 材料を集める ── 入口は裏づき、あらすじだけが残った",
    claim="履歴とクレジットの取得は成立した。残っているのは、ページからあらすじの内容を取り出す部分である。",
    stacks=[
        stack(0, AREA_W, [
            dict(k="lanes", gap=0.16, lanes=[
                [card("履歴の取り込み ── 成立した",
                      "・購入確認メールから 7／7 で公演を特定できた。実運用でも 192 回ぶんが取れた\n"
                      "・記録の日付とメールの日付が 256／258 ＝ 99.2% 一致した\n"
                      "・同じ作品を複数回観た記録を 13／13 で 1 つに束ねられた\n"
                      "・発行元が上位 3 社に集まらなかったので、件名で絞って 1 件ずつ判定する 3 段に変えた",
                      style="br"),
                 card("クレジットの取得 ── 表方は成立し、裏方は捨てた",
                      "・出演 88%・演出 88%・脚本 82% で取れた\n"
                      "・裏方は 16〜42% しか取れず、名簿に足しても AUC が動かなかった（0.500 対 0.661）\n"
                      "・取得率ではなく寄与で判断し、裏方は取りに行かないと決めた",
                      style="br")],
                [card("あらすじ ── ここが最大の未了である",
                      "・規則で切り出すと 37% しか取れず、8% は別の公演の内容を拾った\n"
                      "・人が読めば 58%、公演固有のページに限れば 77% 取れた\n"
                      "・障害は「あらすじが無いこと」ではなく「個別ページに辿り着いていないこと」だった\n"
                      "・LLM 自体の適合率・再現率は、API の枠が別課金のため未測である",
                      style="hi"),
                 card("母集団 ── 被覆が標本で割れた",
                      "・劇場の公式サイトを標本にすると 81% 覆えていた\n"
                      "・当事者が実際に使う経路（ステージナタリーの記事）で測ると 43% だった\n"
                      "・取得しやすさで標本を決めていたので、フィードを取得先に上げる方針にした",
                      style="hi")],
            ]),
            label("この段で命題を書き直した箇所 ── 測って初めて、閾値の置き方が誤っていたと分かった"),
            bullets([
                "裏方の閾値を「取得率 5 割」から「名簿に足したときの寄与」に変えた ── 取れても寄与しないものを追っていた",
                "あらすじの担い手を規則から LLM に変えた ── 取得先を変える手ではなく、担い手を変える手である",
            ]),
            src("検証 004・005・013・014・018・020・022・023・025・027。標本と限界は 1 件ずつ記録に残してある"),
        ], gap=0.10),
    ],
)

# ── 3 枚目 ② 知識を作る ──────────────────────────────────────────
S3 = dict(
    no=3, title="② 知識を作る ── 順位を決めるのは名簿、件数を埋めるのは内容の傾向",
    claim="申告に一致した公演は別枠に出るので、推定は「申告で拾えない公演」で測る。そこでの判別力は 0.751 だった。",
    stacks=[
        stack(0, AREA_W, [
            table(
                ["できる知識", "役割", "分かったこと"],
                [["お気に入り", "無条件に出す",
                  "一致した公演は順位付けの対象から外し、別枠に出す。混ぜると精度が 8.4 ポイント落ちることを"
                  "実測した（0.744 → 0.660）ので、この排他は測って支持された"],
                 ["作り手の名簿", "順位を決める",
                  "申告で拾えない 73 作品での AUC は 0.751。名簿から申告した名前を全部除いても 0.751 で落ちない"
                  "（本人に聞くだけでは作れない知識である）。ただし候補の 92〜95% は一致が 0 件で、週 3〜4 件しか作れない"],
                 ["内容の傾向", "件数を埋める",
                  "候補側に通すと 522／818 件であらすじが取れ、スコアの付く候補が 4.7 倍になった。"
                  "ただし順位は当てていない（AUC 0.551 ± 0.167、B に足して上がった回 0／40）"],
                 ["評価語の表", "記録の画面",
                  "自分の感想の語から極性を作る。順位には使わず、記録を見返す画面に出す"]],
                widths=[0.95, 1.15, 6.96]),
            label("測った結果、企画から外した 2 つ ── 完成後の構成には入らない"),
            bullets([
                "他人のクチコミの読み替え ── 候補の 1.6% にしか付いておらず、109 日の猶予の中に存在しない（検証 029）",
                "会場ごとの当たり率 ── 公演の ◎ は内容への評価で、会場の評価を含まない（検証 002）。客席数は事実として出すだけにした",
            ]),
            label("網 C が空くと何が起きるか ── 3 回目の提示で実際に見えた"),
            bullets([
                "網 B と網 C の両方が理由になった 3 件では 2 件に興味ありが付いたが、網 C だけの 7 件は全滅した",
                "理由が「家族」「時代物」「悲劇」のような広い語 1 つだと、理由として働かない ── 括りの粒度を測る命題に直結する",
            ]),
            src("語の約束 ── **申告**は本人が登録した名前（団体・人・主催・作品・原作者・題材の 22 件）、**推定**は履歴と ◎○△× からシステムが導いた知識である／検証 009・011・012・014・017・018・022・026・029・030・033・034。D を落とした根拠は検証 029、B を中核に移した根拠は検証 030、C の役割分けは検証 033・034、E を降ろした根拠は検証 002 である"),
        ], gap=0.10),
    ],
)

# ── 4 枚目 ③④ 出して、返す ──────────────────────────────────────
S4 = dict(
    no=4, title="③④ 出して、返す ── 反応が推薦を変えるところまで確かめた",
    claim="推薦は 5 回出した。反応を読む処理は成立したが、良くなるかどうかは在庫が尽きて測れていない。",
    stacks=[
        stack(0, AREA_W, [
            label("分かったこと"),
            bullets([
                "候補 818 件のうち基準線を超えたのは 75 件（9%）で、絞り込みそのものは効いている",
                "反応を読む処理を入れると、上位 15 件のうち 8 件が入れ替わった（閾値 3 割に対して成立）"
                " ── 表を作ったことと、読む処理があることは別だと分かった",
                "購入済みの混入は 13% から 0% に、興味なしの再提示は 6 件から 0 件になった。"
                "抑制の鍵を作品単位にして初めて成立した（公演単位ではツアーの別 ID で戻ってくる）",
                "3 回・各 15 件の反応でスコアの判別力を測ると AUC 0.728 で、履歴の中で測った 0.744 とほぼ同じだった",
            ]),
            dict(k="lanes", gap=0.18, lanes=[
                [bars([("1 回目", 57, "興味あり 8／14"),
                       ("2 回目", 40, "6／15"),
                       ("3 回目", 13, "2／15")],
                      baseline=False,
                      note="下がったのは反映が悪いからではなく、スコアの高い順に在庫を"
                           "消費したからである。効果はこの 3 回では判定できない")],
                [card("輪が閉じきらない原因は、画面に反応の入力欄が無いことである",
                      "・いまの反応は会話から入れている\n"
                      "・欄が無いままだと 2 回目の提示に反応が付かず、効果は測れないままになる\n"
                      "・記録を見返す画面は、母集団との比較に作り替えた",
                      style="hi")],
            ]),
            src("検証 015・016・017・019・020・021・023・026。推薦を出した 5 回の記録がそのまま標本になっている"),
        ], gap=0.10),
    ],
)

# ── 5 枚目 残り ──────────────────────────────────────────────────
S5 = dict(
    no=5, title="残っている作業 ── 3 つに絞られた",
    claim="納期内に測るのは 13 項目である。そのうち 3 つが、ほかの項目を止めている。",
    stacks=[
        stack(0, AREA_W, [
            label("止めている 3 つ ── 順序は無く、どれも独立に着手できる"),
            dict(k="lanes", gap=0.14, lanes=[
                [card("あらすじの抽出を LLM に移す",
                      "・上に 4 項目が乗っている（網 C だけの精度・括りの粒度・週あたりの件数・"
                      "「確認できず」の読み）\n"
                      "・偽なら網 C を落とし、企画書を A・B・D の 3 本に書き換える",
                      num="① 材料", style="hi")],
                [card("母集団を増やす",
                      "・3 回の提示で在庫が尽きた。増やさないと、次は広い語 1 つの公演しか出せない\n"
                      "・当事者の主経路で 43% しか覆えていないので、フィードを取得先に上げる",
                      num="① 材料", style="hi")],
                [card("出した理由が役に立ったかを聞く",
                      "・見せた後にしか聞けないので、次の版を出す前に聞き終える\n"
                      "・あわせて提示日を凍結する。今日やらないと、納期後に測る道も閉じる",
                      num="③④ 出力と輪", style="hi")],
            ]),
            label("納期内は測らないもの ── 測らないことと、代わりに何を書くかを決めてある"),
            bullets([
                "P3 の 6 項目は、暫定値または限界として企画書に書く"
                "（例: 推定と申告の比は、履歴の中の代理で 1.34 倍が出ているのでそれを暫定として引く）",
                "凍結の 4 項目は写真からの公演特定が前提で、残り時間では判定が出ない。クチコミ本文を前提にしていた項目は、網 D を落としたので凍結ではなく取り下げに移した",
                "命題の書き方そのものも 6 件直した ── 閾値を「測りやすい量」から決めず、"
                "設計が要求する量か企画書の主張から逆算する",
            ]),
            src("完了条件は、P0〜P2 の全項目に「成立／不成立／部分的に成立」の判定が付き、"
                "スクリプトを再実行して同じ判定が得られることである"),
        ], gap=0.10),
    ],
)

LAYOUTS = [S1, S2, S3, S4, S5]


# ── 出力 ─────────────────────────────────────────────────────────
def draw_pptx(resolved):
    prs = Presentation(str(TEMPLATE))
    ids = prs.slides._sldIdLst
    R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sld in list(ids):
        prs.part.drop_rel(sld.get(R))
        ids.remove(sld)
    for spec in resolved:
        slide = prs.slides.add_slide(prs.slide_layouts[4])  # タイトルのみ
        slide.shapes.title.text = spec["title"]
        for sh in spec["shapes"]:
            draw_shape_pptx(slide, sh)
    prs.save(str(OUT_PPTX))


def draw_html(resolved):
    parts = [HTML_HEAD.replace("たぐり 7 枚版", DECK), '<div class="wrap">',
             f"<h1>{esc(DECK)}</h1>",
             '<p class="lead">検証（#000006）の現在地を、企画書 1 章の流れ（① 材料を集める → '
             '② 知識を作る → ③ 一覧を見る → ④ 評価を付ける）の上に貼った 5 枚である。'
             f'<b>この画面と同じ座標から <code>{OUT_PPTX.name}</code> を生成している。</b>'
             '命題・閾値・標本・限界は <code>docs/verification/</code> の 27 件に残してあり、'
             'ここには「分かったこと」だけを置く。</p>']
    for spec in resolved:
        parts.append(f'<p class="sheetno">SHEET {spec["no"]} ／ {esc(spec["title"])}</p>')
        parts.append('<div class="ab">'
                     f'<div class="ttl">{esc(spec["title"])}</div><div class="rule"></div>'
                     + "".join(html_shape(sh) for sh in spec["shapes"]) + "</div>")
    parts.append("</div>")
    OUT_HTML.write_text("\n".join(parts), encoding="utf-8")


def md_items(items, out, depth=0):
    pad = "  " * depth
    for it in items:
        k = it["k"]
        if k == "label":
            out.append(f"\n{pad}**{it['text']}**\n")
        elif k in ("para", "src"):
            out.append(f"{pad}{it['text']}\n" if k == "para" else f"{pad}※ {it['text']}\n")
        elif k == "card":
            head = f"{it['num']}｜{it['t']}" if it.get("num") else it["t"]
            out.append(f"{pad}- **{head}**")
            for ln in (it.get("b") or "").split("\n"):
                if ln.strip():
                    out.append(f"{pad}  - {ln.lstrip('・')}")
        elif k == "table":
            out.append(f"\n{pad}| " + " | ".join(it["cols"]) + " |")
            out.append(f"{pad}|" + "---|" * len(it["cols"]))
            for r in it["rows"]:
                out.append(f"{pad}| " + " | ".join(r) + " |")
            out.append("")
        elif k == "bars":
            for lb, _v, sub in it["rows"]:
                out.append(f"{pad}- {lb} ── {sub}")
            if it.get("note"):
                out.append(f"{pad}  ※ {it['note']}")
        elif k == "cycle":
            for lb, sub in (it.get("entry") or []):
                out.append(f"{pad}- 〔入口〕**{lb}** ── {sub}")
            for lb, sub in it["nodes"]:
                out.append(f"{pad}- **{lb}** ── {sub}")
            if it.get("interlude"):
                out.append(f"{pad}- 〔分岐〕**{it['interlude'][0]}** ── {it['interlude'][1]}")
        elif k == "lanes":
            for lane in it["lanes"]:
                md_items(lane, out, depth)
        elif k == "group":
            out.append(f"\n{pad}**{it['title']}**")
            md_items(it["items"], out, depth + 1)


def draw_md():
    L = [f"# {DECK}", "",
         "- 対象タスク: [#000006 「たぐり」の前提を実データで検証する]"
         "(../tasks/active/current/000006-verify-taguri-assumptions.md)",
         "- 出力: [000006-verification-status.pptx](000006-verification-status.pptx)"
         " / [プレビュー HTML](000006-verification-status.html)",
         "- 生成: `python3 tools/verify/build_status_slides.py`"
         "（この台本・pptx・プレビューは**同じレイアウト定義から生成される**）", "",
         "検証結果は、企画書 1 章の流れ（① 材料を集める → ② 知識を作る → ③ 一覧を見る → "
         "④ 評価を付ける）のどの段に当たるかで並べてある。**命題・閾値・標本・限界は "
         "`docs/verification/` の 27 件に置き、ここには「分かったこと」だけを書く。**", ""]
    for lay in LAYOUTS:
        L += ["---", "", f"## {lay['no']} 枚目　{lay['title']}", "",
              f"**{lay['claim']}**", ""]
        for st in lay["stacks"]:
            md_items(st["items"], L)
        L.append("")
    OUT_MD.write_text("\n".join(L) + "\n", encoding="utf-8")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="収まるかだけ検査する")
    a = ap.parse_args()
    resolved = [resolve(l) for l in LAYOUTS]
    print(f"収まり検査（本文領域 {AREA_W} x {AREA_H} inch）:")
    for sp in resolved:
        print(f"  {sp['no']} 枚目: 下端 {sp['bottom']:.2f}")
    problems = check(resolved)
    if problems:
        print("\nはみ出し:")
        for pb in problems:
            print("  " + pb)
    else:
        print("すべて枠内に収まっている。")
    if a.check:
        return
    draw_pptx(resolved)
    draw_html(resolved)
    draw_md()
    print(f"\n生成: {OUT_PPTX.name} / {OUT_HTML.name} / {OUT_MD.name}")


if __name__ == "__main__":
    main()
